import os
import json
import glob
import subprocess
import matplotlib.pyplot as plt
import pandas as pd
import seaborn as sns

from tensorboard.backend.event_processing import event_accumulator
from pptx import Presentation
from pptx.util import Inches, Pt
prs = Presentation()
def process_experiment(exp_path):
    
    files = os.listdir(exp_path)
    
    # --- Validate expected files ---
    # TensorBoard log: assume file contains 'tfevents' in its name
    tb_logs = [f for f in files if "tfevents" in f]
    # model
    model_files = [f for f in files if f.lower().endswith(".pth")]
    # CSV files
    csv_files = [f for f in files if f.lower().endswith(".csv")]
    # --- Step 2: Process TensorBoard log to plot loss curves ---
    tb_log_file = os.path.join(exp_path, tb_logs[0])
    try:
        ea = event_accumulator.EventAccumulator(tb_log_file)
        ea.Reload()
        # Extract scalars for training and validation loss
        train_loss = ea.Scalars('Loss/Train')
        val_loss = ea.Scalars('Loss/Validation')
        train_acc = ea.Scalars('Accuracy/Train')
        val_acc = ea.Scalars('Accuracy/Validation')
        params = ['trial_number', 'initial_lr', 'weight_decay', 'momentum', 'swa_start', 'swa_lr']
        metrics = ['accuracy', 'auroc', 'aupr_success', 'aupr', 'fpr', 'tnr', 'aurc', 'eaurc', 'augrc']
        check = True
        try:
            trial_number = ea.Scalars(f'final_Metrics/train_accuracy')
        except KeyError:
            check = False  # or set a default value or log a warning
            return -1
        if check:
            params_list = []
            for param in params:
                params_list.append(ea.Scalars(f'Params/{param}'))
            
            
            train_metrics_final = []
            for metric in metrics:
                train_metrics_final.append(ea.Scalars(f'final_Metrics/train_{metric}'))
            val_metrics_final = []
            for metric in metrics:
                val_metrics_final.append(ea.Scalars(f'final_Metrics/val_{metric}'))
            test_metrics_final = []
            for metric in metrics:
                test_metrics_final.append(ea.Scalars(f'final_Metrics/test_{metric}'))
            train_metrics_final = [sublist[0].value for sublist in train_metrics_final]
            val_metrics_final = [sublist[0].value for sublist in val_metrics_final]
            test_metrics_final = [sublist[0].value for sublist in test_metrics_final]
            params_list = [sublist[0].value for sublist in params_list]

    except Exception as e:
        print(f"Error reading TensorBoard log in {exp_path}: {e}")
        return

    # Plot train and validation loss on the same plot
    plt.figure(figsize=(8, 6))
    train_steps = [x.step for x in train_loss]
    train_values = [x.value for x in train_loss]
    val_steps = [x.step for x in val_loss]
    val_values = [x.value for x in val_loss]
    FINAL_TRAIN_LOSS = train_values[-1]
    FINAL_VAL_LOSS = val_values[-1]
    plt.plot(train_steps, train_values, label="Train Loss")
    plt.plot(val_steps, val_values, label="Validation Loss")
    plt.xlabel("Epoch")
    plt.ylabel("Loss")
    plt.title("Loss Curve")
    plt.legend()
    loss_plot_path = os.path.join(exp_path, "loss_plot.png")
    plt.savefig(loss_plot_path)
    plt.close()
    
    plt.figure(figsize=(8, 6))
    train_steps = [x.step for x in train_acc]
    train_values = [x.value for x in train_acc]
    val_steps = [x.step for x in val_acc]
    val_values = [x.value for x in val_acc]
    plt.plot(train_steps, train_values, label="Train Accuracy")
    plt.plot(val_steps, val_values, label="Validation Accuracy")
    plt.xlabel("Step")
    plt.ylabel("Accuracy")
    plt.title("Accuracy Plot")
    plt.legend()
    acc_plot_path = os.path.join(exp_path, "accuracy_plot.png")
    plt.savefig(acc_plot_path)
    plt.close()
    
    # --- Step 3: Process CSV files to create confidence plots ---
    confidence_plots = {} # key: csv file name, value: plot file path
    for csv_file in csv_files:
        csv_path = os.path.join(exp_path, csv_file)
        try:
            df = pd.read_csv(csv_path)
        except Exception as e:
            print(f"Error reading CSV file {csv_path}: {e}")
            continue

        # Expected columns: 'sample name', 'confidence value', 'prediction', 'ground truth'
        expected_cols = {'Confidence (%)', 'correct'}
        if not expected_cols.issubset(df.columns):
            print(f"CSV file {csv_file} does not have the expected columns.")
            continue

        plt.figure(figsize=(8, 6))
        # Plot histograms for confidence: green for correct, red for wrong
        plt.hist(df[df['correct']]['Confidence (%)'], bins=50, color='green',
                 alpha=0.6, density=True, label='Success')
        plt.hist(df[~df['correct']]['Confidence (%)'], bins=50, color='red',
                 alpha=0.6, density=True, label='Error')
        plt.xlim([0, 100])
        plt.xlabel("Confidence")
        plt.ylabel("Density")
        plt.title(f"Confidence Distribution - {csv_file}")
        plt.legend()
        plot_name = csv_file.replace(".csv", "_confidence.png")
        plot_path = os.path.join(exp_path, plot_name)
        plt.savefig(plot_path)
        plt.close()
        confidence_plots[csv_file] = plot_path
    
    # --- Step 4: Read metrics from tensorboard log and create a PowerPoint slide ---
    # Create a single slide PowerPoint presentation with the metrics as title text and images in the slide.
    
    # Use a blank slide layout (index may vary based on template)
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    aka = ''
    # Add a title textbox for metrics
    for jj in range(0, len(metrics) - 1, 2):
        aka = aka + (f'{metrics[jj]}: {val_metrics_final[jj]:.2f}\t')
        aka = aka + (f'{metrics[jj+1]}: {val_metrics_final[jj+1]:.2f}\n')
    title_text = f"validation metrics:\n{aka}"
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(10)
    aka = ''
    # Add a title textbox for metrics
    '''
    for jj in range(len(params)):
        if params[jj] != 'trial_number':
            aka = aka + (f'{params[jj]}: {params_list[jj]:.5f}\n')
        else:
            aka = aka + (f'{params[jj]}: {int(params_list[jj])}\n')
    '''
    title_text = f"params:\n{aka}"
    title_box = slide.shapes.add_textbox(Inches(7.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    for paragraph in title_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(10)
    # Collect all generated images (loss plot + confidence plots)
    images = [loss_plot_path, acc_plot_path] 
    # Arrange images in a grid (up to 4 images)
    positions = [
        (Inches(0.2), Inches(4.5)),
        (Inches(5), Inches(4.5)),
        (Inches(0.2), Inches(1.5)),
        (Inches(5), Inches(1.5))
    ]
    for i, img in enumerate(images):
        if i >= len(positions):
            break
        left, top = positions[i]
        slide.shapes.add_picture(img, left, top, width=Inches(5), height=Inches(3))    
    csv_files = [i for i in csv_files if ('test' not in i) and ('DELETE' not in i) ]
    print(csv_files)
    for i, csv_file in enumerate(csv_files):
        left, top = positions[i+2]
        slide.shapes.add_picture(confidence_plots[csv_file], left, top, width=Inches(5), height=Inches(3))    
    
    # --- Mark experiment as processed ---
    with open(os.path.join(exp_path, ".processed"), 'w') as f:
        f.write("processed")
    print(f"Processed experiment at {exp_path}")
    
    experiment_data = {'experiment': os.path.basename(exp_path)}
    for i, metric in enumerate(metrics):
        experiment_data[metric] = val_metrics_final[i]
    experiment_data['train_loss@200'] = FINAL_TRAIN_LOSS
    experiment_data['val_loss@200'] = FINAL_VAL_LOSS
    experiment_data['loss_diff'] = FINAL_VAL_LOSS - FINAL_TRAIN_LOSS 
    print(experiment_data)
    return experiment_data
def main():  
    experiments_root = '/home/amax/machairas/FMFP-edge-idid/hailo_src/haht_augrc_resnet_18/'
    
    if not os.path.exists(experiments_root):
        print("No EXPERIMENTS directory found.")
        return
    
    all_experiments_data = []
    # Process each experiment folder that is new (i.e. does not contain a '.processed' marker)
    directories = [ d for d in os.listdir(experiments_root) if os.path.isdir(os.path.join(experiments_root, d))]

    for exp_folder in sorted(directories, key=int)[:-1]:
        if exp_folder != '21':
            continue
        exp_path = os.path.join(experiments_root, exp_folder)
        #if ".processed" not in os.listdir(exp_path):
        print(f"Processing {exp_path}...")
        experiment_metrics  = process_experiment(exp_path)
        if experiment_metrics == -1:
            continue
        all_experiments_data.append(experiment_metrics)
    pptx_path = os.path.join(experiments_root, "experiment_report_final.pptx")
    prs.save(pptx_path)
    # --- Create CSV file ---
    if all_experiments_data:
        df = pd.DataFrame(all_experiments_data)
        csv_path = os.path.join(experiments_root, "metrics_summary_final.csv")
        df.to_csv(csv_path, index=False)
        print(f"CSV summary saved to {csv_path}")
    else:
        print("No experiment data to save.")
if __name__ == "__main__":
    main()